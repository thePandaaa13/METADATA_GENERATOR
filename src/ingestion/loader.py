import json 
from pathlib import Path
from typing import List,Dict

from src.config import SAMPLE_ASSETS_DIR, ASSET_CONTEXT_FILE

def load_asset_context()->Dict[str,dict]:
    with open(ASSET_CONTEXT_FILE, "r" ) as f:
        context_list = json.load(f)
    
    context_by_filename = {item["filename"]: item for item in context_list}
    return context_by_filename

    
def load_all_assets() -> List[dict]:
    context_by_filename = load_asset_context()
    assets = []

    supported_extensions = ["*.txt", "*.pdf", "*.docx", "*.pptx"]

    for pattern in supported_extensions:
        for file_path in SAMPLE_ASSETS_DIR.glob(pattern):
            content = extract_text(file_path)
            filename = file_path.name
            context = context_by_filename.get(filename, {})

            assets.append({
                "filename": filename,
                "content": content,
                "context": context,
            })

    return assets

if __name__ == "__main__":
    assets = load_all_assets()
    print(f"Total assets loaded : {len(assets)}\n")

    for asset in assets:
        print(f" File : {asset['filename']}")
        print(f"Context : {asset['context']}")
        print(f"content preview: {asset['content'][:80]}")
        print("-----")

  
# loading the brochures
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_path) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_path)->str:
    doc = Document(file_path)
    text ="/n".join(paragraph.text for paragraph in doc.paragraphs)
    return text

def extract_text_from_pptx(file_path) ->str:
    prs = Presentation(file_path)
    text=""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame :
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs : 
                        text += run.text + " "
        text += "\n"
    return text                

def extract_text(file_path) -> str:
    
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    elif suffix == ".docx":
        return extract_text_from_docx(file_path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    elif suffix == ".txt":
        with open(file_path, "r") as f:
            return f.read()
    else:
        raise ValueError(f"Unsupported file type: {suffix}")